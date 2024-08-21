"""
Author: Carlos L. Gray
Date: August 2024

Description:
This script provides functionality for reading and indexing various document types into a persistent Chroma vector store. 
The script includes the following key components:
1. Initialization of embeddings using Hugging Face models and the Chroma vector store with persistence enabled.
2. Functions to read content from different file formats, including text, PDF, DOCX, XLSX, and PPTX files.
3. A function to check if a document has already been indexed in the Chroma store to avoid duplication.
4. An indexing function that processes and chunks the document content, then indexes the chunks into the Chroma vector store.
5. The text chunking is done using LangChain's `RecursiveCharacterTextSplitter` to ensure efficient storage and retrieval.

This script is particularly useful for organizations that want to efficiently index and search large volumes of document data, 
leveraging AI-powered vector stores for rapid retrieval and contextual understanding of the content.

Dependencies:
- langchain_chroma: Integrates the Chroma vector store for efficient document storage and retrieval.
- langchain_huggingface: Provides Hugging Face embeddings for converting text into vector representations.
- PyMuPDF (fitz): Used for reading content from PDF files.
- docx: A module to read Microsoft Word documents.
- openpyxl: A module to read Excel files.
- python-pptx: A module to read PowerPoint files.
- Python 3.x
"""

import os
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Initialize the embeddings and Chroma database with persistence
persist_directory = "./chroma_db"  # Specify the directory to persist the Chroma database
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
chroma_document_store = Chroma(embedding_function=embeddings, persist_directory=persist_directory)

#Read in the chunking values, if not set use defaults
chunk_size = os.getenv('CHUNK_SIZE', 1000)
chunk_overlap = os.getenv('CHUNK_OVERLAP', 200)
max_batch_size = os.getenv('MAX_BATCH_SIZE', 5461)

def read_file_content(file_path):
    """Read content from a file and return it as a string."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.txt':
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            print(f"UnicodeDecodeError with utf-8 on {file_path}, trying cp1252...")
            try:
                with open(file_path, 'r', encoding='cp1252') as f:
                    return f.read()
            except UnicodeDecodeError:
                print(f"UnicodeDecodeError with cp1252 on {file_path}, trying 'ignore' errors...")
                with open(file_path, 'r', encoding='cp1252', errors='ignore') as f:
                    return f.read()
    elif ext == '.pdf':
        text = ""
        try:
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
            return text
        except Exception as e:  # Catch all exceptions related to PDF parsing
            print(f"Skipping {file_path} due to PDF error: {e}")
            return ""
    elif ext == '.docx':
        try:
            doc = docx.Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs])
        except Exception as e:
            print(f"Error reading {file_path}: {e}")
            return ""
    elif ext == '.xlsx':
        wb = openpyxl.load_workbook(file_path)
        text = ""
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text += ' '.join([str(cell) for cell in row if cell]) + '\n'
        return text
    elif ext == '.pptx':
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return ""

def is_document_indexed(file_name, chroma_document_store):
    """Check if a document is already indexed in the Chroma store."""
    collection = chroma_document_store._client.get_or_create_collection(chroma_document_store._collection.name)
    documents = collection.get(include=["metadatas", "documents"])
    
    for meta in documents["metadatas"]:
        if meta.get("name") == file_name:
            return True
    return False

def index_documents_from_file_server(directory_path, chunk_size=chunk_size, chunk_overlap=chunk_overlap, max_batch_size=max_batch_size):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    for root, dirs, files in os.walk(directory_path):
        for file in files:
            file_path = os.path.join(root, file)
            if is_document_indexed(file, chroma_document_store):
                print(f"Document {file} is already indexed. Skipping...")
                continue

            print(f"Processing file: {file_path}")  # Debugging output
            content = read_file_content(file_path)
            if content:
                print(f"Indexing document: {file}")  # Debugging output
                texts = text_splitter.split_text(content)

                # Process texts in batches to avoid exceeding the maximum batch size
                for i in range(0, len(texts), max_batch_size):
                    batch_texts = texts[i:i + max_batch_size]
                    chroma_document_store.add_texts(batch_texts, metadatas=[{"name": file}] * len(batch_texts))

            else:
                print(f"Skipping empty or failed document: {file}")  # Debugging output

    print(f"Documents indexed into collection: {chroma_document_store._collection.name}")
    
index_documents_from_file_server(os.getenv('DOCUMENT_DIRECTORY'), chunk_size=chunk_size, chunk_overlap=chunk_overlap)


